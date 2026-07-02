"""PPTX: each figure on its own slide (18 slides total)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

ROOT = Path('/home/hyunjin/CLAUDE/Pd_DMC/research-pd-dmc')
FIG = ROOT / 'reports/predft_advisor_figures'
OUT = ROOT / 'reports/DMC_Pd_advisor_figures.pptx'

p = Presentation()
p.slide_width = Inches(13.33); p.slide_height = Inches(7.5)
blank = p.slide_layouts[6]

from PIL import Image
def add_slide(png_path, title):
    s = p.slides.add_slide(blank)
    box = s.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.5))
    tf = box.text_frame; tf.text = title
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(22); r.font.bold = True
    # fit fully into 12.5×6.3 box (aspect-preserving)
    im = Image.open(png_path)
    aspect = im.width / im.height
    box_w, box_h = 12.5, 6.3
    if aspect > box_w/box_h:
        w, h = box_w, box_w/aspect
    else:
        w, h = box_h*aspect, box_h
    x = (13.33 - w) / 2
    y = 0.9 + (box_h - h) / 2
    s.shapes.add_picture(str(png_path), Inches(x), Inches(y),
                          width=Inches(w), height=Inches(h))

SLABS = [('S1','S1: Pd(100) (Pd⁰)'),
         ('S2','S2: 1ML PdO(101)/Pd(100) (Pd⁰+Pd²⁺)'),
         ('S3','S3: O-rich PdO(100) (Pd²⁺)'),
         ('S3b','S3b: Pd-rich PdO(100) (Pd²⁺)'),
         ('S4','S4: PdO₂(110) (Pd⁴⁺)')]
for sid, ttl in SLABS:
    add_slide(FIG/'slab_sideviews_final'/f'{sid}_side.png', f'Slab side view — {ttl}')

VIOLINS = [('F13a_CO.png','E_bind distribution — CO*'),
           ('F13b_CH3O.png','E_bind distribution — CH₃O*'),
           ('F13c_coads.png','E_bind distribution — co-ads (CO* + CH₃O*)')]
for fn, ttl in VIOLINS:
    add_slide(FIG/'violin_individual'/fn, ttl)

ADS = [('S1','Pd(100), Pd⁰'),('S2','1ML PdO(101)/Pd(100)'),
       ('S3','O-rich PdO(100)'),('S3b','Pd-rich PdO(100)'),
       ('S4','PdO₂(110), Pd⁴⁺')]
for ads_kind, ads_lbl in [('CO','CO*'),('CH3O','CH₃O*')]:
    for sid, name in ADS:
        add_slide(FIG/'top1_adsorbate_sideviews_v2'/f'{sid}_{ads_kind}_top1.png',
                  f'Top-1 {ads_lbl} — {sid}: {name}')

p.save(str(OUT))
print(f'Saved: {OUT} ({len(p.slides)} slides)')
