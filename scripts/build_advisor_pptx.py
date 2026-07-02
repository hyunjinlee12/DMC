"""Build PPTX: slabs + violin plots + top-1 adsorbates, dpi 300 figures."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

ROOT = Path('/home/hyunjin/CLAUDE/Pd_DMC/research-pd-dmc')
FIG = ROOT / 'reports/predft_advisor_figures'
OUT = ROOT / 'reports/DMC_Pd_advisor_figures.pptx'

p = Presentation()
p.slide_width = Inches(13.33); p.slide_height = Inches(7.5)
blank = p.slide_layouts[6]

def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.5))
    tf = box.text_frame; tf.text = text
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.LEFT
    for r in p0.runs: r.font.size = Pt(22); r.font.bold = True

# Slide 1: 5 slabs side view
s = p.slides.add_slide(blank); add_title(s, 'Converged slab side views (2×2×1, bc-plane)')
imgs = ['S1_side.png','S2_side.png','S3_side.png','S3b_side.png','S4_side.png']
labels = ['S1: Pd(100) (Pd⁰)','S2: 1ML PdO(101)/Pd(100)','S3: O-rich PdO(100)',
          'S3b: Pd-rich PdO(100)','S4: PdO₂(110) (Pd⁴⁺)']
w = 2.5
for i, (f, lab) in enumerate(zip(imgs, labels)):
    x = Inches(0.2 + i*2.6); y = Inches(1.0)
    s.shapes.add_picture(str(FIG/'slab_sideviews_final'/f), x, y, width=Inches(w))
    box = s.shapes.add_textbox(x, Inches(6.2), Inches(w), Inches(0.4))
    tf = box.text_frame; tf.text = lab
    for r in tf.paragraphs[0].runs: r.font.size = Pt(10); r.font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 2: Violin plot (a) CO*
s = p.slides.add_slide(blank); add_title(s, 'E_bind distribution — CO*')
s.shapes.add_picture(str(FIG/'violin_individual/F13a_CO.png'),
                    Inches(2.5), Inches(0.8), height=Inches(6.5))

# Slide 3: Violin plot (b) CH3O*
s = p.slides.add_slide(blank); add_title(s, 'E_bind distribution — CH₃O*')
s.shapes.add_picture(str(FIG/'violin_individual/F13b_CH3O.png'),
                    Inches(2.5), Inches(0.8), height=Inches(6.5))

# Slide 4: Violin plot (c) co-ads
s = p.slides.add_slide(blank); add_title(s, 'E_bind distribution — co-ads (CO* + CH₃O*)')
s.shapes.add_picture(str(FIG/'violin_individual/F13c_coads.png'),
                    Inches(2.5), Inches(0.8), height=Inches(6.5))

# Slide 5: Top-1 CO
s = p.slides.add_slide(blank); add_title(s, 'Top-1 lowest-E CO* (intramol-valid)')
s.shapes.add_picture(str(FIG/'top1_adsorbate_sideviews_v2/ALL_CO_top1.png'),
                    Inches(0.3), Inches(1.0), width=Inches(12.7))

# Slide 6: Top-1 CH3O
s = p.slides.add_slide(blank); add_title(s, 'Top-1 lowest-E CH₃O* (intramol-valid)')
s.shapes.add_picture(str(FIG/'top1_adsorbate_sideviews_v2/ALL_CH3O_top1.png'),
                    Inches(0.3), Inches(1.0), width=Inches(12.7))

p.save(str(OUT))
print(f'Saved: {OUT}')
